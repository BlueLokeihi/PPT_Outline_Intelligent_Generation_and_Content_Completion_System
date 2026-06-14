from __future__ import annotations

import io
import json
import re
import shutil
import time
import uuid
from dataclasses import asdict
from pathlib import Path
from typing import Any, Dict, List, Literal, Optional

from dotenv import load_dotenv
from fastapi import FastAPI, File, Form, UploadFile
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import FileResponse, JSONResponse
from pydantic import BaseModel, Field

from llm.client import build_provider, load_models_config
from monitoring import monitor, record_api_call, record_api_call_async, user_facing_error
from outline.evaluation import compute_quality
from outline.generator import generate_once
from outline.io_utils import maybe_truncate
from outline.prompt_strategies import PromptOptions, build_messages
from rag.embedder import load_embed_config
from rag.enricher import (
    confidence_from_coverage,
    enrich_page,
    snippets_from_research_chunks,
)
from rag.query_rewriter import page_context_from_outline
from rag.research import research_page
from rag.retriever import HybridRetriever
from rag.store import corpus_path

# 加载 .env 文件中的环境变量
PROJECT_ROOT = Path(__file__).resolve().parents[1]
BACKEND_ROOT = PROJECT_ROOT / "backend"
MODELS_CONFIG = BACKEND_ROOT / "config" / "models.json"
ENV_FILE = BACKEND_ROOT / ".env"
ALLOWED_CORPUS_EXTS = {".pdf", ".txt", ".md"}

# 自动加载 .env 文件
if ENV_FILE.exists():
    load_dotenv(ENV_FILE)
    print(f"[OK] Loaded env: {ENV_FILE}")
else:
    print(f"[WARN] .env not found: {ENV_FILE}")


class MessageItem(BaseModel):
    role: Literal["system", "user", "assistant"]
    text: str = ""


class OutlineRequest(BaseModel):
    conversationId: str | None = None
    messages: List[MessageItem] = Field(default_factory=list)
    pdfText: str = ""
    provider: str = "qwen"
    strategy: Literal["baseline", "few_shot", "cot_silent"] = "baseline"
    schema: Literal["on", "off"] = "on"
    minSlides: int = Field(default=10, ge=1, le=100)
    maxSlides: int = Field(default=18, ge=1, le=100)
    # —— RAG 相关（默认关闭，旧客户端无影响） ——
    useRag: bool = False
    corpusId: str = ""
    ragMode: Literal["vector", "bm25", "hybrid"] = "hybrid"
    ragRewriteThreshold: float = Field(default=0.45, ge=0.0, le=1.0)
    ragMaxRounds: int = Field(default=2, ge=1, le=3)
    ragMaxSnippets: int = Field(default=8, ge=1, le=20)
    ragSkipLowConfidence: bool = False
    ragLowThreshold: float = Field(default=0.4, ge=0.0, le=1.0)
    # 并发处理页数；默认 3，太高容易触发 provider QPS 限速
    ragConcurrency: int = Field(default=3, ge=1, le=8)
    # 额外上下文
    fileSources: List[FileSourceItem] = Field(default_factory=list)
    webResults: List[WebResultItem] = Field(default_factory=list)


class OutlineSavePage(BaseModel):
    title: str = ""
    bullets: List[str] = Field(default_factory=list)
    notes: str = ""
    evidences: Optional[List[Dict[str, Any]]] = None
    conflicts: Optional[List[Dict[str, Any]]] = None


class OutlineSaveChapter(BaseModel):
    title: str = ""
    pages: List[OutlineSavePage] = Field(default_factory=list)


class OutlineSaveBody(BaseModel):
    title: str = ""
    assumptions: Optional[List[str]] = None
    chapters: List[OutlineSaveChapter] = Field(default_factory=list)


class OutlineSaveRequest(BaseModel):
    conversationId: Optional[str] = None
    outline: OutlineSaveBody


class OutlineVersionRequest(BaseModel):
    conversationId: Optional[str] = None
    outline: OutlineSaveBody
    sourceType: Literal["generated", "edited", "rag", "restored"] = "edited"
    provider: Optional[str] = None
    strategy: Optional[str] = None
    schemaMode: Optional[Literal["on", "off"]] = None
    useRag: bool = False
    summary: str = ""


class OutlineExportRequest(BaseModel):
    conversationId: Optional[str] = None
    outline: OutlineSaveBody
    format: Literal["markdown", "html", "pptx", "json"] = "markdown"
    report: Optional[Dict[str, Any]] = None


class CorpusBuildRequest(BaseModel):
    corpusId: str
    provider: str = "qwen"
    model: str = ""
    chunkSize: int = Field(default=500, ge=100, le=2000)
    overlap: int = Field(default=80, ge=0, le=500)


class CorpusRenameRequest(BaseModel):
    corpusId: str
    newCorpusId: str


class QuestionnaireRequest(BaseModel):
    topic: str
    provider: str = "qwen"


class WebSearchRequest(BaseModel):
    query: str
    maxResults: int = Field(default=8, ge=1, le=20)


class FileSourceItem(BaseModel):
    name: str = ""
    text: str = ""


class WebResultItem(BaseModel):
    title: str = ""
    url: str = ""
    snippet: str = ""


_QUESTIONNAIRE_SYSTEM = (
    "你是PPT大纲生成助手的前置问卷模块。\n\n"
    "用户已通过表单提供了场景、受众、目标、风格、深度等基本参数。"
    "你的任务是：在此基础上，针对该主题进一步挖掘内容层面的细节，"
    "提出2-3个能显著影响大纲结构与内容深度的选择题。\n\n"
    "提问原则：\n"
    "- 不要重复表单已有的字段（场景/受众/风格/深度等）\n"
    "- 聚焦于内容本身：核心论点侧重、数据与案例的使用方式、叙事结构、需要重点攻克的难点等\n"
    "- 每题给出3个有实质区别的选项（15字以内），全部使用中文\n"
    "- 必须始终输出问题，needs_questionnaire 始终为 true\n\n"
    "严格返回JSON，不含其他文字：\n"
    '{"needs_questionnaire":true,"questions":['
    '{"id":"q1","question":"...","options":[{"id":"a","label":"..."},{"id":"b","label":"..."},{"id":"c","label":"..."}],'
    '"allow_custom":true,"allow_ai_decide":true}]}'
)


app = FastAPI(title="PPT Outline Local API", version="0.1.0")
app.add_middleware(
    CORSMiddleware,
    allow_origins=["http://localhost:5173", "http://127.0.0.1:5173"],
    allow_origin_regex=r"http://(localhost|127\.0\.0\.1)(:\d+)?",
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)


def compose_topic_text(
    messages: List[MessageItem],
    pdf_text: str,
    file_sources: Optional[List[Dict[str, str]]] = None,
    web_results: Optional[List[Dict[str, str]]] = None,
) -> str:
    lines: List[str] = ["你将基于以下对话与资料生成 PPT 大纲。"]

    if pdf_text.strip():
        lines.append("\n[PDF资料摘录]\n")
        lines.append(pdf_text.strip())

    if file_sources:
        for fs in file_sources:
            name = fs.get("name", "附件")
            text = (fs.get("text") or "").strip()
            if text:
                lines.append(f"\n[参考资料：{name}]\n")
                lines.append(text[:8000])

    if web_results:
        lines.append("\n[网络搜索结果]\n")
        for i, r in enumerate(web_results[:8], 1):
            title = r.get("title", "")
            url = r.get("url", "")
            snippet = r.get("snippet", "")
            lines.append(f"{i}. 【{title}】")
            lines.append(f"   来源：{url}")
            lines.append(f"   摘要：{snippet}")

    lines.append("\n[多轮对话]\n")
    for item in messages:
        text = item.text.strip()
        if not text:
            continue
        lines.append(f"{item.role}: {text}")

    return "\n".join(lines).strip()


def run_outline(request: OutlineRequest) -> Dict[str, Any]:
    provider = request.provider.strip() or "qwen"
    strategy = request.strategy
    schema = request.schema
    min_slides = request.minSlides
    max_slides = request.maxSlides

    if min_slides > max_slides:
        min_slides, max_slides = max_slides, min_slides

    topic_text = compose_topic_text(
        messages=request.messages,
        pdf_text=request.pdfText,
        file_sources=[{"name": fs.name, "text": fs.text} for fs in request.fileSources],
        web_results=[{"title": r.title, "url": r.url, "snippet": r.snippet} for r in request.webResults],
    )
    topic_text = maybe_truncate(topic_text, max_chars=24000)

    client, model, defaults, _provider_cfg = build_provider(provider, config_path=str(MODELS_CONFIG))

    opts = PromptOptions(
        strategy=strategy,
        enforce_schema=(schema == "on"),
        min_slides=min_slides,
        max_slides=max_slides,
    )
    messages = build_messages(topic_text, opts)

    result = generate_once(
        client,
        provider=provider,
        strategy=strategy,
        enforce_schema=(schema == "on"),
        messages=messages,
        model=model,
        temperature=defaults.temperature,
        top_p=defaults.top_p,
        max_tokens=defaults.max_tokens,
    )

    response: Dict[str, Any] = {
        "ok": bool(result.ok),
        "provider": provider,
        "strategy": strategy,
        "schema": schema,
        "elapsedS": round(float(result.elapsed_s), 3),
    }

    if result.ok and result.outline is not None:
        response["outline"] = result.outline
        response["quality"] = asdict(
            compute_quality(result.outline, min_slides=min_slides, max_slides=max_slides)
        )

        if request.useRag and request.corpusId.strip():
            try:
                rag_meta = _run_rag_pipeline(request, response["outline"])
                response["rag"] = rag_meta
                response["elapsedS"] = round(
                    response["elapsedS"] + float(rag_meta.get("elapsed_s") or 0.0), 3
                )
            except Exception as exc:  # noqa: BLE001
                response["rag"] = {"used": False, "error": str(exc)}
        try:
            version = _write_outline_version(
                conversation_id=request.conversationId,
                outline=response["outline"],
                source_type="rag" if response.get("rag", {}).get("used") else "generated",
                provider=provider,
                strategy=strategy,
                schema_mode=schema,
                use_rag=bool(response.get("rag", {}).get("used")),
            )
            response["version"] = {
                "versionId": version["versionId"],
                "createdAt": version["createdAt"],
                "sourceType": version["sourceType"],
                "summary": version["summary"],
            }
        except Exception:
            pass
        return response

    response["error"] = result.error or "生成失败"
    return response


def _clean_text(value: str, fallback: str) -> str:
    text = (value or "").strip()
    return text if text else fallback


def _clean_list(items: Optional[List[str]]) -> List[str]:
    if not items:
        return []
    return [item.strip() for item in items if item and item.strip()]


def _normalize_outline(outline: OutlineSaveBody) -> Dict[str, Any]:
    chapters: List[Dict[str, Any]] = []
    for chapter in outline.chapters:
        pages: List[Dict[str, Any]] = []
        for page in chapter.pages:
            page_payload: Dict[str, Any] = {
                "title": _clean_text(page.title, "未命名页面"),
                "bullets": _clean_list(page.bullets),
                "notes": page.notes or "",
            }
            if page.evidences:
                page_payload["evidences"] = page.evidences
            if page.conflicts:
                page_payload["conflicts"] = page.conflicts
            pages.append(page_payload)
        chapters.append({
            "title": _clean_text(chapter.title, "未命名章节"),
            "pages": pages,
        })

    payload: Dict[str, Any] = {
        "title": _clean_text(outline.title, "未命名大纲"),
        "chapters": chapters,
    }

    if outline.assumptions:
        payload["assumptions"] = _clean_list(outline.assumptions)

    return payload


def _safe_token(value: Optional[str]) -> str:
    if not value:
        return "session"
    token = re.sub(r"[^a-zA-Z0-9_-]+", "_", value).strip("_")
    return token[:48] or "session"


def _invalidate_corpus_index(corpus_id: str) -> bool:
    index_dir = BACKEND_ROOT / "rag" / "indexes" / corpus_id
    if not index_dir.exists():
        return False
    shutil.rmtree(index_dir, ignore_errors=True)
    return True


def _scan_corpus_dir(root: Path) -> tuple[int, int, float]:
    file_count = 0
    total_bytes = 0
    latest_mtime = 0.0
    for p in root.rglob("*"):
        if not p.is_file():
            continue
        stat = p.stat()
        file_count += 1
        total_bytes += stat.st_size
        latest_mtime = max(latest_mtime, stat.st_mtime)
    return file_count, total_bytes, latest_mtime


def _format_ts(ts: float) -> str:
    if ts <= 0:
        return ""
    return time.strftime("%Y-%m-%d %H:%M:%S", time.localtime(ts))


def _versions_root(conversation_id: Optional[str]) -> Path:
    return BACKEND_ROOT / "outline" / "versions" / _safe_token(conversation_id)


def _outline_summary(outline: Dict[str, Any]) -> str:
    chapters = outline.get("chapters") or []
    page_count = sum(len(ch.get("pages") or []) for ch in chapters)
    return f"{len(chapters)} chapters, {page_count} pages"


def _write_outline_version(
    *,
    conversation_id: Optional[str],
    outline: Dict[str, Any],
    source_type: str,
    provider: Optional[str] = None,
    strategy: Optional[str] = None,
    schema_mode: Optional[str] = None,
    use_rag: bool = False,
    summary: str = "",
) -> Dict[str, Any]:
    version_id = f"{time.strftime('%Y%m%d_%H%M%S')}_{uuid.uuid4().hex[:8]}"
    created_at = time.strftime("%Y-%m-%dT%H:%M:%S")
    payload = {
        "versionId": version_id,
        "conversationId": conversation_id or "",
        "createdAt": created_at,
        "sourceType": source_type,
        "provider": provider,
        "strategy": strategy,
        "schema": schema_mode,
        "useRag": use_rag,
        "summary": summary or _outline_summary(outline),
        "outline": outline,
    }
    out_dir = _versions_root(conversation_id)
    out_dir.mkdir(parents=True, exist_ok=True)
    (out_dir / f"{version_id}.json").write_text(
        json.dumps(payload, ensure_ascii=False, indent=2),
        encoding="utf-8",
    )
    return payload


def _read_version_file(version_id: str) -> Dict[str, Any]:
    base = BACKEND_ROOT / "outline" / "versions"
    for path in base.glob(f"*/{version_id}.json"):
        return json.loads(path.read_text(encoding="utf-8"))
    raise FileNotFoundError(f"version not found: {version_id}")


def _build_markdown(outline: Dict[str, Any], report: Optional[Dict[str, Any]] = None) -> str:
    lines = [f"# {outline.get('title') or 'PPT Outline'}"]
    assumptions = outline.get("assumptions") or []
    if assumptions:
        lines.extend(["", "## 假设前提"])
        lines.extend([f"- {item}" for item in assumptions])
    for c_idx, chapter in enumerate(outline.get("chapters") or [], 1):
        lines.extend(["", f"## {c_idx}. {chapter.get('title') or '未命名章节'}"])
        for p_idx, page in enumerate(chapter.get("pages") or [], 1):
            lines.extend(["", f"### {c_idx}.{p_idx} {page.get('title') or '未命名页面'}"])
            for bullet in page.get("bullets") or []:
                lines.append(f"- {bullet}")
            if page.get("notes"):
                lines.extend(["", f"> {page['notes']}"])
            evidences = page.get("evidences") or []
            if evidences:
                lines.append("")
                lines.append("来源：")
                for ev in evidences:
                    source = ev.get("source") or "unknown"
                    chunk = ev.get("chunk_index")
                    suffix = f" #chunk{chunk}" if isinstance(chunk, int) else ""
                    lines.append(f"- {source}{suffix}: {ev.get('text', '')[:160]}")
            conflicts = page.get("conflicts") or []
            if conflicts:
                lines.append("")
                lines.append("冲突提示：")
                for item in conflicts:
                    lines.append(f"- {item.get('message') or item}")
    if report:
        lines.extend(["", "## 验收报告", "```json", json.dumps(report, ensure_ascii=False, indent=2), "```"])
    return "\n".join(lines) + "\n"


def _build_html(outline: Dict[str, Any], report: Optional[Dict[str, Any]] = None) -> str:
    def esc(value: Any) -> str:
        return (
            str(value)
            .replace("&", "&amp;")
            .replace("<", "&lt;")
            .replace(">", "&gt;")
            .replace('"', "&quot;")
        )

    body = [f"<h1>{esc(outline.get('title') or 'PPT Outline')}</h1>"]
    for c_idx, chapter in enumerate(outline.get("chapters") or [], 1):
        body.append(f"<section><h2>{c_idx}. {esc(chapter.get('title') or '未命名章节')}</h2>")
        for p_idx, page in enumerate(chapter.get("pages") or [], 1):
            body.append(f"<article><h3>{c_idx}.{p_idx} {esc(page.get('title') or '未命名页面')}</h3><ul>")
            for bullet in page.get("bullets") or []:
                body.append(f"<li>{esc(bullet)}</li>")
            body.append("</ul>")
            if page.get("notes"):
                body.append(f"<p class='notes'>{esc(page['notes'])}</p>")
            evidences = page.get("evidences") or []
            if evidences:
                body.append("<details><summary>Sources</summary><ol>")
                for ev in evidences:
                    body.append(f"<li><strong>{esc(ev.get('source', 'unknown'))}</strong>: {esc(ev.get('text', ''))}</li>")
                body.append("</ol></details>")
            conflicts = page.get("conflicts") or []
            if conflicts:
                body.append("<div class='conflict'><strong>Conflict notes</strong><ul>")
                for item in conflicts:
                    body.append(f"<li>{esc(item.get('message') if isinstance(item, dict) else item)}</li>")
                body.append("</ul></div>")
            body.append("</article>")
        body.append("</section>")
    if report:
        body.append(f"<section><h2>Acceptance Report</h2><pre>{esc(json.dumps(report, ensure_ascii=False, indent=2))}</pre></section>")

    return """<!doctype html>
<html lang="zh-CN">
<head>
  <meta charset="utf-8" />
  <title>PPT Outline Export</title>
  <style>
    body { font-family: -apple-system, BlinkMacSystemFont, "Segoe UI", sans-serif; margin: 40px; color: #172033; line-height: 1.55; }
    article { border-top: 1px solid #dbeafe; padding: 14px 0; }
    h1 { color: #1d4ed8; }
    h2 { color: #0f766e; }
    .notes { background: #f8fafc; border-left: 4px solid #60a5fa; padding: 10px 12px; }
    .conflict { background: #fff7ed; border: 1px solid #fdba74; border-radius: 8px; padding: 10px 12px; }
    pre { background: #0f172a; color: #e2e8f0; padding: 16px; border-radius: 8px; overflow: auto; }
  </style>
</head>
<body>
""" + "\n".join(body) + "\n</body>\n</html>\n"


def _build_pptx(outline: Dict[str, Any], out_path: Path) -> None:
    try:
        from pptx import Presentation
        from pptx.util import Pt
    except Exception as exc:  # noqa: BLE001
        raise RuntimeError("PPTX export requires python-pptx. Please install backend requirements.") from exc

    prs = Presentation()
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = outline.get("title") or "PPT Outline"
    title_slide.placeholders[1].text = "Generated by PPT Outline System"

    for c_idx, chapter in enumerate(outline.get("chapters") or [], 1):
        section = prs.slides.add_slide(prs.slide_layouts[1])
        section.shapes.title.text = f"{c_idx}. {chapter.get('title') or '未命名章节'}"
        section.placeholders[1].text = f"{len(chapter.get('pages') or [])} pages"
        for p_idx, page in enumerate(chapter.get("pages") or [], 1):
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = f"{c_idx}.{p_idx} {page.get('title') or '未命名页面'}"
            frame = slide.placeholders[1].text_frame
            frame.clear()
            for b_idx, bullet in enumerate(page.get("bullets") or ["（暂无要点）"]):
                para = frame.paragraphs[0] if b_idx == 0 else frame.add_paragraph()
                para.text = str(bullet)
                para.level = 0
                para.font.size = Pt(18)
            notes = page.get("notes") or ""
            if notes:
                slide.notes_slide.notes_text_frame.text = notes
    prs.save(out_path)


def _detect_conflicts(evidences: List[Dict[str, Any]], confidence: str) -> List[Dict[str, Any]]:
    """Lightweight, explainable conflict flagging for reviewer-facing RAG output.

    This is intentionally conservative: it marks "possible conflicts" when retrieved
    evidence comes from multiple sources and contains different numeric claims.
    """
    if len(evidences) < 2:
        return []
    numbers_by_source: Dict[str, set[str]] = {}
    for ev in evidences:
        source = str(ev.get("source") or "unknown")
        text = str(ev.get("text") or "")
        nums = set(re.findall(r"\d+(?:\.\d+)?%?|\d{4}", text))
        if nums:
            numbers_by_source.setdefault(source, set()).update(nums)
    if len(numbers_by_source) < 2:
        return []
    all_sets = list(numbers_by_source.values())
    shared = set.intersection(*all_sets) if all_sets else set()
    union = set.union(*all_sets) if all_sets else set()
    if len(union - shared) < 2 and confidence != "low":
        return []
    return [{
        "type": "possible_numeric_conflict",
        "severity": "medium" if confidence != "low" else "high",
        "message": "Multiple sources contain different numeric claims; please confirm before presenting.",
        "sources": sorted(numbers_by_source.keys()),
        "signals": sorted(union - shared)[:12],
    }]


# ---- RAG 管道（research + enrich，全在内存中流转） ------------------------

def _process_rag_page(
    *,
    request: OutlineRequest,
    outline: Dict[str, Any],
    chapter: Dict[str, Any],
    c_idx: int,
    p_idx: int,
    retriever: HybridRetriever,
    embed_cfg: Optional[Any],
    llm_client: Any,
    llm_model: str,
    overall_topic: str,
) -> Dict[str, Any]:
    """处理单页：research → (可选 skip) → enrich。

    返回一个 dict：
      {
        "c_idx": ..., "p_idx": ...,
        "page_meta": {...},                 # 写入 rag.page_meta 的条目
        "page_updates": Optional[{...}],    # 命中成功时要写回 page 的字段
      }
    page_updates 为 None 表示不要改写原 page；否则其中包含 bullets/notes/evidences/conflicts。
    """
    pages = chapter.get("pages") or []
    page = pages[p_idx]
    page_title = str(page.get("title", ""))
    ctx = page_context_from_outline(outline, c_idx, p_idx)

    # ---- research：query 改写 + 多轮检索 ----
    r = research_page(
        ctx,
        retriever=retriever,
        embed_cfg=embed_cfg,
        llm_client=llm_client,
        llm_model=llm_model,
        llm_provider=request.provider,
        mode=request.ragMode,
        k=6,
        recall_k=20,
        rrf_k=60,
        rewrite_threshold=request.ragRewriteThreshold,
        max_rounds=request.ragMaxRounds,
        verbose=False,
    )
    if r.get("error"):
        return {
            "c_idx": c_idx, "p_idx": p_idx,
            "page_meta": {
                "chapter_idx": c_idx, "page_idx": p_idx,
                "page_title": page_title,
                "enrichment": "research_failed",
                "error": r["error"],
            },
            "page_updates": None,
        }

    rounds = r.get("rounds") or []
    coverage = float(rounds[-1].get("coverage") or 0.0) if rounds else 0.0
    confidence = confidence_from_coverage(coverage)

    # ---- 低置信度 skip ----
    if request.ragSkipLowConfidence and coverage < request.ragLowThreshold:
        return {
            "c_idx": c_idx, "p_idx": p_idx,
            "page_meta": {
                "chapter_idx": c_idx, "page_idx": p_idx,
                "page_title": page_title,
                "enrichment": "skipped_low_confidence",
                "coverage": round(coverage, 4),
                "confidence": confidence,
            },
            "page_updates": None,
        }

    # ---- enrich：浓缩为 bullets+notes+evidences ----
    snippets = snippets_from_research_chunks(r.get("final_chunks") or [])
    er = enrich_page(
        overall_topic=overall_topic,
        chapter_title=str(chapter.get("title", "")),
        page_title=page_title,
        original_bullets=ctx.bullets,
        original_notes=ctx.notes,
        snippets=snippets,
        coverage=coverage,
        client=llm_client,
        model=llm_model,
        provider=request.provider,
        temperature=0.3,
        max_snippets=request.ragMaxSnippets,
    )
    if er.error or not er.bullets:
        return {
            "c_idx": c_idx, "p_idx": p_idx,
            "page_meta": {
                "chapter_idx": c_idx, "page_idx": p_idx,
                "page_title": page_title,
                "enrichment": "failed",
                "coverage": round(coverage, 4),
                "confidence": confidence,
                "error": er.error,
            },
            "page_updates": None,
        }

    updates: Dict[str, Any] = {
        "bullets": er.bullets,
        "notes": er.notes,
    }
    conflicts: List[Dict[str, Any]] = []
    if er.evidences:
        updates["evidences"] = er.evidences
        conflicts = _detect_conflicts(er.evidences, confidence)
        if conflicts:
            updates["conflicts"] = conflicts

    return {
        "c_idx": c_idx, "p_idx": p_idx,
        "page_meta": {
            "chapter_idx": c_idx, "page_idx": p_idx,
            "page_title": page_title,
            "enrichment": "ok",
            "coverage": round(coverage, 4),
            "confidence": confidence,
            "used_source_ids": er.used_source_ids,
            "n_evidences": len(er.evidences),
            "conflicts": conflicts,
        },
        "page_updates": updates,
    }


def _run_rag_pipeline(
    request: OutlineRequest, outline: Dict[str, Any]
) -> Dict[str, Any]:
    """跑 research + enrich（多页并发），写回 evidences/bullets/notes。

    并发度由 request.ragConcurrency 控制（默认 3）。同一 page 的 research+enrich 仍是
    顺序的（共一个 worker），但不同 page 并发；这样既缩短 wall time，也保证每页失败可
    独立归因。
    """
    from concurrent.futures import ThreadPoolExecutor, as_completed

    t_start = time.time()
    corpus_id = request.corpusId.strip()
    out_dir = corpus_path(corpus_id)

    retriever = HybridRetriever(out_dir)
    manifest = retriever.store.load_manifest()    # 不存在则抛 FileNotFoundError

    embed_cfg = None
    if request.ragMode in ("vector", "hybrid"):
        embed_cfg = load_embed_config(
            str(MODELS_CONFIG),
            provider=request.provider,
            model=manifest.get("embedding_model"),
        )

    llm_client, llm_model, _defaults, _cfg = build_provider(
        request.provider, config_path=str(MODELS_CONFIG)
    )

    overall_topic = str(outline.get("title", ""))
    chapters = outline.get("chapters") or []

    # —— 提交所有 page 任务到线程池 ——
    tasks: List[tuple[int, int, Dict[str, Any]]] = []
    for c_idx, chapter in enumerate(chapters):
        for p_idx, _page in enumerate(chapter.get("pages") or []):
            tasks.append((c_idx, p_idx, chapter))

    concurrency = max(1, min(int(request.ragConcurrency), 8))
    results_by_key: Dict[tuple[int, int], Dict[str, Any]] = {}

    if not tasks:
        return {
            "used": True,
            "corpus": corpus_id,
            "mode": request.ragMode,
            "embedding_model": manifest.get("embedding_model"),
            "index_size": manifest.get("size"),
            "rewrite_threshold": request.ragRewriteThreshold,
            "max_rounds": request.ragMaxRounds,
            "concurrency": concurrency,
            "page_meta": [],
            "elapsed_s": round(time.time() - t_start, 3),
        }

    with ThreadPoolExecutor(max_workers=concurrency) as pool:
        futures = {
            pool.submit(
                _process_rag_page,
                request=request,
                outline=outline,
                chapter=chapter,
                c_idx=c_idx,
                p_idx=p_idx,
                retriever=retriever,
                embed_cfg=embed_cfg,
                llm_client=llm_client,
                llm_model=llm_model,
                overall_topic=overall_topic,
            ): (c_idx, p_idx)
            for c_idx, p_idx, chapter in tasks
        }
        for fut in as_completed(futures):
            key = futures[fut]
            try:
                results_by_key[key] = fut.result()
            except Exception as exc:  # noqa: BLE001
                c_idx, p_idx = key
                results_by_key[key] = {
                    "c_idx": c_idx, "p_idx": p_idx,
                    "page_meta": {
                        "chapter_idx": c_idx, "page_idx": p_idx,
                        "page_title": "",
                        "enrichment": "exception",
                        "error": str(exc),
                    },
                    "page_updates": None,
                }

    # —— 按 (c_idx, p_idx) 升序写回 outline + 组装 page_meta ——
    page_meta: List[Dict[str, Any]] = []
    for c_idx, chapter in enumerate(chapters):
        for p_idx, page in enumerate(chapter.get("pages") or []):
            entry = results_by_key.get((c_idx, p_idx))
            if entry is None:
                continue
            updates = entry.get("page_updates")
            if updates is not None:
                page["bullets"] = updates["bullets"]
                page["notes"] = updates["notes"]
                if "evidences" in updates:
                    page["evidences"] = updates["evidences"]
                elif "evidences" in page:
                    del page["evidences"]
                if "conflicts" in updates:
                    page["conflicts"] = updates["conflicts"]
                elif "conflicts" in page:
                    del page["conflicts"]
            page_meta.append(entry["page_meta"])

    return {
        "used": True,
        "corpus": corpus_id,
        "mode": request.ragMode,
        "embedding_model": manifest.get("embedding_model"),
        "index_size": manifest.get("size"),
        "rewrite_threshold": request.ragRewriteThreshold,
        "max_rounds": request.ragMaxRounds,
        "concurrency": concurrency,
        "page_meta": page_meta,
        "elapsed_s": round(time.time() - t_start, 3),
    }


@app.post("/api/ping")
def api_ping() -> Dict[str, Any]:
    return {"ok": True, "message": "http api is ready"}


@app.get("/api/monitoring/usage")
def api_monitoring_usage() -> Dict[str, Any]:
    return {"ok": True, "metrics": monitor.snapshot()}


@app.post("/api/monitoring/reset")
def api_monitoring_reset() -> Dict[str, Any]:
    return {"ok": True, "metrics": monitor.reset()}


@app.post("/api/questionnaire")
async def api_questionnaire(req: QuestionnaireRequest) -> Dict[str, Any]:
    """Generate clarifying questions for a PPT topic using LLM."""
    try:
        provider = req.provider.strip() or "qwen"
        client, model, defaults, _ = build_provider(provider, config_path=str(MODELS_CONFIG))

        # 直接 await client.chat() — 避免在 uvicorn event loop 中嵌套 asyncio.run()
        resp = await record_api_call_async(
            api_type="llm",
            provider=provider,
            operation="questionnaire",
            func=lambda: client.chat(
                model=model,
                messages=[
                    {"role": "system", "content": _QUESTIONNAIRE_SYSTEM},
                    {"role": "user", "content": f"用户PPT需求：{req.topic}"},
                ],
                temperature=0.3,
                top_p=defaults.top_p,
                max_tokens=800,
            ),
        )
        raw = str(resp.choices[0].message.content or "").strip()

        # Extract JSON from response
        json_match = re.search(r"\{.*\}", raw, re.DOTALL)
        if json_match:
            data = json.loads(json_match.group())
            questions = data.get("questions") or []
            # 强制触发：只要 LLM 生成了有效问题就弹出问卷，不依赖 LLM 自报 needs_questionnaire
            if questions:
                return {"ok": True, "needs_questionnaire": True, "questions": questions}
        return {"ok": True, "needs_questionnaire": False, "questions": []}
    except Exception as exc:  # noqa: BLE001
        return {"ok": True, "needs_questionnaire": False, "questions": [], "error": user_facing_error(exc, api_type="llm")}


def _search_jina(query: str, max_results: int) -> List[Dict[str, str]]:
    """Jina Search (s.jina.ai)。主引擎，无需 Key 亦可用，配置 JINA_API_KEY 额度更高。"""
    import os
    import urllib.parse
    import requests as _req

    encoded = urllib.parse.quote(query, safe="")
    url = f"https://s.jina.ai/{encoded}"
    headers: Dict[str, str] = {"Accept": "application/json", "X-Retain-Images": "none"}
    api_key = os.getenv("JINA_API_KEY", "").strip()
    if api_key:
        headers["Authorization"] = f"Bearer {api_key}"

    resp = _req.get(url, headers=headers, timeout=20)
    resp.raise_for_status()
    data = resp.json()

    results: List[Dict[str, str]] = []
    for item in (data.get("data") or [])[:max_results]:
        snippet = (item.get("description") or item.get("content") or "").strip()
        results.append({
            "title": (item.get("title") or "").strip(),
            "url":   (item.get("url")   or "").strip(),
            "snippet": snippet[:240] + ("…" if len(snippet) > 240 else ""),
        })
    return results


def _search_serp(query: str, max_results: int) -> List[Dict[str, str]]:
    """SerpAPI (Google)。兜底引擎，Free Plan 250次/月。"""
    import os
    import requests as _req

    api_key = os.getenv("SERP_API_KEY", "").strip()
    if not api_key:
        raise RuntimeError("SERP_API_KEY 未配置")

    params = {
        "q": query,
        "api_key": api_key,
        "engine": "google",
        "num": max_results,
        "hl": "zh-cn",
        "gl": "cn",
    }
    resp = _req.get("https://serpapi.com/search.json", params=params, timeout=20)
    resp.raise_for_status()
    data = resp.json()

    results: List[Dict[str, str]] = []
    for item in (data.get("organic_results") or [])[:max_results]:
        snippet = (item.get("snippet") or "").strip()
        results.append({
            "title": (item.get("title") or "").strip(),
            "url":   (item.get("link")  or "").strip(),
            "snippet": snippet[:240] + ("…" if len(snippet) > 240 else ""),
        })
    return results


@app.post("/api/search/web")
async def api_web_search(req: WebSearchRequest) -> Dict[str, Any]:
    """双引擎搜索：Jina 优先，失败时自动切换 SerpAPI 兜底。"""
    import asyncio as _asyncio

    def _search() -> Dict[str, Any]:
        # ── 主引擎：Jina ──
        try:
            results = record_api_call(
                api_type="web_search",
                provider="jina",
                operation="search",
                func=lambda: _search_jina(req.query, req.maxResults),
            )
            if results:
                return {"ok": True, "results": results, "engine": "jina"}
        except Exception as jina_exc:  # noqa: BLE001
            jina_err = user_facing_error(jina_exc, api_type="web_search")
        else:
            jina_err = "Jina 返回空结果"

        # ── 兜底：SerpAPI ──
        serp_err = "SerpAPI 返回空结果"
        try:
            results = record_api_call(
                api_type="web_search",
                provider="serp",
                operation="search",
                func=lambda: _search_serp(req.query, req.maxResults),
            )
            if results:
                return {"ok": True, "results": results, "engine": "serp"}
        except Exception as serp_exc:  # noqa: BLE001
            serp_err = user_facing_error(serp_exc, api_type="web_search")

        return {
            "ok": False,
            "results": [],
            "error": "两个搜索引擎均不可用，请稍后重试",
            "jina_error": jina_err,
            "serp_error": serp_err,
        }

    try:
        loop = _asyncio.get_event_loop()
        result = await loop.run_in_executor(None, _search)
        return result
    except Exception:  # noqa: BLE001
        return {"ok": False, "results": [], "error": "搜索服务暂时不可用，请稍后重试或手动上传文件"}


@app.post("/api/extract/docx")
async def api_extract_docx(file: UploadFile = File(...)) -> Dict[str, Any]:
    """提取 Word (.docx) 文件文本。"""
    try:
        import docx  # python-docx
        content = await file.read()
        doc = docx.Document(io.BytesIO(content))
        paragraphs = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
        # Also extract table text
        for table in doc.tables:
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                if cells:
                    paragraphs.append(" | ".join(cells))
        text = "\n".join(paragraphs)
        return {"ok": True, "text": text, "fileName": file.filename, "charCount": len(text)}
    except Exception as exc:  # noqa: BLE001
        return {"ok": False, "error": str(exc)}


@app.post("/api/extract/pptx")
async def api_extract_pptx(file: UploadFile = File(...)) -> Dict[str, Any]:
    """提取 PowerPoint (.pptx) 文件文本。"""
    try:
        from pptx import Presentation
        content = await file.read()
        prs = Presentation(io.BytesIO(content))
        texts: List[str] = []
        for slide_idx, slide in enumerate(prs.slides, 1):
            slide_texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_texts.append(shape.text.strip())
            if slide_texts:
                texts.append(f"[幻灯片 {slide_idx}]\n" + "\n".join(slide_texts))
        text = "\n\n".join(texts)
        return {"ok": True, "text": text, "fileName": file.filename, "charCount": len(text)}
    except Exception as exc:  # noqa: BLE001
        return {"ok": False, "error": str(exc)}


@app.get("/api/runtime-info")
def api_runtime_info() -> Dict[str, Any]:
    try:
        cfg = load_models_config(str(MODELS_CONFIG))
        providers = list((cfg.get("providers") or {}).keys())
        return {
            "ok": True,
            "providers": providers,
            "strategies": ["baseline", "few_shot", "cot_silent"],
        }
    except Exception:
        return {
            "ok": False,
            "providers": ["qwen", "glm", "deepseek"],
            "strategies": ["baseline", "few_shot", "cot_silent"],
        }


@app.get("/api/rag/corpora")
def api_rag_corpora() -> Dict[str, Any]:
    """列出已建好的索引（backend/rag/indexes/<id>/）。"""
    base = BACKEND_ROOT / "rag" / "indexes"
    if not base.exists():
        return {"ok": True, "corpora": []}
    items: List[Dict[str, Any]] = []
    for d in sorted(base.iterdir()):
        if not d.is_dir():
            continue
        manifest_path = d / "manifest.json"
        if not manifest_path.exists():
            continue
        try:
            import json as _json
            manifest = _json.loads(manifest_path.read_text(encoding="utf-8"))
            items.append({
                "id": d.name,
                "size": manifest.get("size"),
                "dim": manifest.get("dim"),
                "embedding_model": manifest.get("embedding_model"),
                "built_at": manifest.get("built_at"),
                "has_bm25": (d / "tokens.jsonl").exists(),
            })
        except Exception:
            continue
    return {"ok": True, "corpora": items}


@app.get("/api/rag/corpus-sources")
def api_rag_corpus_sources() -> Dict[str, Any]:
    """列出 corpus 目录中的知识库来源（含是否已建索引）。"""
    base = BACKEND_ROOT / "rag" / "corpus"
    if not base.exists():
        return {"ok": True, "sources": []}
    items: List[Dict[str, Any]] = []
    for d in sorted(base.iterdir()):
        if not d.is_dir():
            continue
        file_count, total_bytes, latest_mtime = _scan_corpus_dir(d)
        index_dir = BACKEND_ROOT / "rag" / "indexes" / d.name
        manifest_path = index_dir / "manifest.json"
        manifest: Dict[str, Any] = {}
        has_index = manifest_path.exists()
        if has_index:
            try:
                manifest = json.loads(manifest_path.read_text(encoding="utf-8"))
            except Exception:
                manifest = {}
        items.append({
            "id": d.name,
            "file_count": file_count,
            "total_bytes": total_bytes,
            "source_updated_at": _format_ts(latest_mtime),
            "has_index": has_index,
            "size": manifest.get("size") if has_index else None,
            "dim": manifest.get("dim") if has_index else None,
            "embedding_model": manifest.get("embedding_model") if has_index else "",
            "built_at": manifest.get("built_at") if has_index else "",
            "has_bm25": (index_dir / "tokens.jsonl").exists() if has_index else False,
        })
    return {"ok": True, "sources": items}


@app.post("/api/rag/corpora/build")
def api_rag_corpus_build(request: CorpusBuildRequest) -> Dict[str, Any]:
    try:
        corpus_id = _safe_token(request.corpusId)
        src_dir = BACKEND_ROOT / "rag" / "corpus" / corpus_id
        if not src_dir.exists():
            return {"ok": False, "error": f"corpus directory not found: {src_dir}"}
        from rag.index import main as build_index
        args = [
            "--corpus", corpus_id,
            "--provider", request.provider,
            "--chunk-size", str(request.chunkSize),
            "--overlap", str(request.overlap),
        ]
        if request.model.strip():
            args.extend(["--model", request.model.strip()])
        code = build_index(args)
        return {"ok": code == 0, "corpusId": corpus_id, "exitCode": code}
    except Exception as exc:  # noqa: BLE001
        return {"ok": False, "error": str(exc)}


@app.post("/api/rag/corpora/rename")
def api_rag_corpus_rename(request: CorpusRenameRequest) -> Dict[str, Any]:
    try:
        src_raw = request.corpusId.strip()
        dst_raw = request.newCorpusId.strip()
        if not src_raw or not dst_raw:
            return {"ok": False, "error": "corpusId and newCorpusId are required"}
        safe_src = _safe_token(src_raw)
        safe_dst = _safe_token(dst_raw)
        if safe_dst != dst_raw:
            return {"ok": False, "error": "newCorpusId must use letters, numbers, underscore, or hyphen"}
        if safe_src == safe_dst:
            return {"ok": False, "error": "newCorpusId must be different"}

        src_dir = BACKEND_ROOT / "rag" / "corpus" / safe_src
        dst_dir = BACKEND_ROOT / "rag" / "corpus" / safe_dst
        if not src_dir.exists():
            return {"ok": False, "error": f"corpus not found: {safe_src}"}
        if dst_dir.exists():
            return {"ok": False, "error": f"corpus already exists: {safe_dst}"}

        src_index = BACKEND_ROOT / "rag" / "indexes" / safe_src
        dst_index = BACKEND_ROOT / "rag" / "indexes" / safe_dst
        if dst_index.exists():
            return {"ok": False, "error": f"index already exists for target: {safe_dst}"}

        shutil.move(str(src_dir), str(dst_dir))
        if src_index.exists():
            shutil.move(str(src_index), str(dst_index))

        return {"ok": True, "corpusId": safe_src, "newCorpusId": safe_dst}
    except Exception as exc:  # noqa: BLE001
        return {"ok": False, "error": str(exc)}


@app.post("/api/rag/upload")
async def api_rag_upload(
    corpusId: str = Form(...),
    files: List[UploadFile] = File(...),
) -> Dict[str, Any]:
    """上传文件到语料库目录（不自动建索引）。"""
    try:
        corpus_id = _safe_token(corpusId)
        if not corpus_id:
            return {"ok": False, "error": "corpusId is required"}
        invalid: List[str] = []
        for f in files:
            filename = Path(f.filename or "").name
            if not filename:
                invalid.append("(empty)")
                continue
            if Path(filename).suffix.lower() not in ALLOWED_CORPUS_EXTS:
                invalid.append(filename)
        if invalid:
            allowed = ", ".join(sorted(ALLOWED_CORPUS_EXTS))
            return {"ok": False, "error": f"不支持的文件类型: {', '.join(invalid)}（仅支持 {allowed}）"}
        dest_dir = BACKEND_ROOT / "rag" / "corpus" / corpus_id
        dest_dir.mkdir(parents=True, exist_ok=True)
        saved: List[str] = []
        for f in files:
            filename = Path(f.filename or "upload").name
            if not filename:
                continue
            dest = dest_dir / filename
            content = await f.read()
            dest.write_bytes(content)
            saved.append(filename)
        invalidated = _invalidate_corpus_index(corpus_id)
        return {
            "ok": True,
            "corpusId": corpus_id,
            "saved": saved,
            "count": len(saved),
            "indexInvalidated": invalidated,
        }
    except Exception as exc:  # noqa: BLE001
        return {"ok": False, "error": str(exc)}


@app.get("/api/rag/corpora/{corpus_id}/files")
def api_rag_corpus_files(corpus_id: str) -> Dict[str, Any]:
    """列出语料库目录中的文件。"""
    try:
        src_dir = BACKEND_ROOT / "rag" / "corpus" / _safe_token(corpus_id)
        if not src_dir.exists():
            return {"ok": True, "files": []}
        files = [
            {"name": p.name, "size": p.stat().st_size}
            for p in sorted(src_dir.iterdir())
            if p.is_file()
        ]
        return {"ok": True, "files": files}
    except Exception as exc:  # noqa: BLE001
        return {"ok": False, "error": str(exc)}


@app.delete("/api/rag/corpora/{corpus_id}/files/{filename}")
def api_rag_corpus_file_delete(corpus_id: str, filename: str) -> Dict[str, Any]:
    """删除语料库目录中的单个文件。"""
    try:
        safe_id = _safe_token(corpus_id)
        safe_name = Path(filename).name
        path = BACKEND_ROOT / "rag" / "corpus" / safe_id / safe_name
        if not path.exists():
            return {"ok": False, "error": "file not found"}
        path.unlink()
        invalidated = _invalidate_corpus_index(safe_id)
        return {"ok": True, "indexInvalidated": invalidated}
    except Exception as exc:  # noqa: BLE001
        return {"ok": False, "error": str(exc)}


@app.delete("/api/rag/corpora/{corpus_id}")
def api_rag_corpus_delete(corpus_id: str) -> Dict[str, Any]:
    """删除整个语料库（含索引）。"""
    try:
        safe_id = _safe_token(corpus_id)
        corpus_dir = BACKEND_ROOT / "rag" / "corpus" / safe_id
        index_dir = BACKEND_ROOT / "rag" / "indexes" / safe_id
        if not corpus_dir.exists() and not index_dir.exists():
            return {"ok": False, "error": "corpus not found"}
        if corpus_dir.exists():
            shutil.rmtree(corpus_dir, ignore_errors=True)
        if index_dir.exists():
            shutil.rmtree(index_dir, ignore_errors=True)
        return {"ok": True}
    except Exception as exc:  # noqa: BLE001
        return {"ok": False, "error": str(exc)}


@app.post("/api/outline")
def api_outline(request: OutlineRequest) -> Dict[str, Any]:
    try:
        return run_outline(request)
    except Exception as exc:  # noqa: BLE001
        return {"ok": False, "error": str(exc)}


@app.post("/api/outline/save")
def api_outline_save(request: OutlineSaveRequest) -> Dict[str, Any]:
    try:
        outline_payload = _normalize_outline(request.outline)
        out_dir = BACKEND_ROOT / "outline" / "output"
        out_dir.mkdir(parents=True, exist_ok=True)
        timestamp = time.strftime("%Y%m%d_%H%M%S")
        token = _safe_token(request.conversationId)
        filename = f"{timestamp}_edited_{token}.json"
        out_path = out_dir / filename
        out_path.write_text(
            json.dumps(outline_payload, ensure_ascii=False, indent=2),
            encoding="utf-8",
        )
        return {
            "ok": True,
            "file": filename,
            "relativePath": f"outline/output/{filename}",
        }
    except Exception as exc:  # noqa: BLE001
        return {"ok": False, "error": str(exc)}


@app.post("/api/outline/versions")
def api_outline_version_create(request: OutlineVersionRequest) -> Dict[str, Any]:
    try:
        outline_payload = _normalize_outline(request.outline)
        version = _write_outline_version(
            conversation_id=request.conversationId,
            outline=outline_payload,
            source_type=request.sourceType,
            provider=request.provider,
            strategy=request.strategy,
            schema_mode=request.schemaMode,
            use_rag=request.useRag,
            summary=request.summary,
        )
        return {"ok": True, "version": {k: version.get(k) for k in (
            "versionId", "conversationId", "createdAt", "sourceType",
            "provider", "strategy", "schema", "useRag", "summary",
        )}}
    except Exception as exc:  # noqa: BLE001
        return {"ok": False, "error": str(exc)}


@app.get("/api/outline/versions")
def api_outline_versions(conversationId: Optional[str] = None) -> Dict[str, Any]:
    try:
        out_dir = _versions_root(conversationId)
        versions: List[Dict[str, Any]] = []
        if out_dir.exists():
            for path in sorted(out_dir.glob("*.json"), reverse=True):
                data = json.loads(path.read_text(encoding="utf-8"))
                versions.append({k: data.get(k) for k in (
                    "versionId", "conversationId", "createdAt", "sourceType",
                    "provider", "strategy", "schema", "useRag", "summary",
                )})
        return {"ok": True, "versions": versions}
    except Exception as exc:  # noqa: BLE001
        return {"ok": False, "error": str(exc), "versions": []}


@app.get("/api/outline/versions/{version_id}")
def api_outline_version_get(version_id: str) -> Dict[str, Any]:
    try:
        return {"ok": True, "version": _read_version_file(version_id)}
    except Exception as exc:  # noqa: BLE001
        return {"ok": False, "error": str(exc)}


@app.post("/api/outline/versions/{version_id}/restore")
def api_outline_version_restore(version_id: str) -> Dict[str, Any]:
    try:
        version = _read_version_file(version_id)
        restored = _write_outline_version(
            conversation_id=version.get("conversationId"),
            outline=version["outline"],
            source_type="restored",
            provider=version.get("provider"),
            strategy=version.get("strategy"),
            schema_mode=version.get("schema"),
            use_rag=bool(version.get("useRag")),
            summary=f"Restored from {version_id}",
        )
        return {
            "ok": True,
            "outline": version["outline"],
            "version": {k: restored.get(k) for k in (
                "versionId", "conversationId", "createdAt", "sourceType",
                "provider", "strategy", "schema", "useRag", "summary",
            )},
        }
    except Exception as exc:  # noqa: BLE001
        return {"ok": False, "error": str(exc)}


@app.post("/api/outline/export")
def api_outline_export(request: OutlineExportRequest):
    try:
        outline_payload = _normalize_outline(request.outline)
        out_dir = BACKEND_ROOT / "outline" / "exports"
        out_dir.mkdir(parents=True, exist_ok=True)
        timestamp = time.strftime("%Y%m%d_%H%M%S")
        token = _safe_token(request.conversationId)
        base = f"{timestamp}_{token}"

        if request.format == "json":
            path = out_dir / f"{base}.json"
            path.write_text(
                json.dumps({"outline": outline_payload, "report": request.report}, ensure_ascii=False, indent=2),
                encoding="utf-8",
            )
            return FileResponse(path, media_type="application/json", filename=path.name)

        if request.format == "html":
            path = out_dir / f"{base}.html"
            path.write_text(_build_html(outline_payload, request.report), encoding="utf-8")
            return FileResponse(path, media_type="text/html; charset=utf-8", filename=path.name)

        if request.format == "pptx":
            path = out_dir / f"{base}.pptx"
            _build_pptx(outline_payload, path)
            return FileResponse(
                path,
                media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                filename=path.name,
            )

        path = out_dir / f"{base}.md"
        path.write_text(_build_markdown(outline_payload, request.report), encoding="utf-8")
        return FileResponse(path, media_type="text/markdown; charset=utf-8", filename=path.name)
    except Exception as exc:  # noqa: BLE001
        return JSONResponse(status_code=500, content={"ok": False, "error": str(exc)})


if __name__ == "__main__":
    import uvicorn

    uvicorn.run("http_server:app", host="127.0.0.1", port=8000, reload=False)

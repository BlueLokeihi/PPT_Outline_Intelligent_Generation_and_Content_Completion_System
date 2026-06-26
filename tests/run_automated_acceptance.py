from __future__ import annotations

import argparse
import io
import json
import os
import subprocess
import sys
import time
import uuid
from datetime import datetime
from pathlib import Path
from typing import Any, Callable


PROJECT_ROOT = Path(__file__).resolve().parents[1]
BACKEND_ROOT = PROJECT_ROOT / "backend"
RESULTS_DIR = PROJECT_ROOT / "test_results"


SAMPLE_OUTLINE = {
    "title": "人工智能发展简史",
    "assumptions": ["测试用静态大纲，用于版本、导出、Schema 和质量指标验证"],
    "chapters": [
        {
            "title": "人工智能的起源",
            "pages": [
                {
                    "title": "图灵测试与学科诞生",
                    "bullets": ["图灵提出机器智能判定问题", "达特茅斯会议提出人工智能概念", "早期研究聚焦符号推理"],
                    "notes": "说明 AI 概念形成的历史背景。",
                },
                {
                    "title": "早期成果与第一次寒冬",
                    "bullets": ["逻辑理论家展示自动推理潜力", "感知机推动神经网络雏形", "算力和算法限制导致预期落差"],
                    "notes": "强调技术能力与社会预期之间的差距。",
                },
            ],
        },
        {
            "title": "从专家系统到机器学习",
            "pages": [
                {
                    "title": "专家系统的商业价值",
                    "bullets": ["专家系统编码领域规则", "XCON 展示企业应用价值", "维护成本限制长期扩展"],
                    "notes": "解释专家系统为何先繁荣后衰退。",
                },
                {
                    "title": "统计学习方法兴起",
                    "bullets": ["SVM 等方法提升分类性能", "Deep Blue 成为标志性事件", "数据驱动方法逐步成为主流"],
                    "notes": "连接符号主义向机器学习范式的转变。",
                },
            ],
        },
        {
            "title": "深度学习与大模型时代",
            "pages": [
                {
                    "title": "深度学习突破",
                    "bullets": ["AlexNet 引爆图像识别突破", "AlphaGo 展示强化学习能力", "Transformer 改变自然语言处理"],
                    "notes": "展示深度学习的关键里程碑。",
                },
                {
                    "title": "生成式 AI 的发展方向",
                    "bullets": ["ChatGPT 推动大众化应用", "多模态能力持续增强", "未来关注可信、安全和高效部署"],
                    "notes": "总结当代 AI 的应用价值与挑战。",
                    "evidences": [
                        {
                            "text": "ChatGPT released in 2022 and accelerated consumer adoption of generative AI.",
                            "source": "test_fixture",
                            "score": 0.91,
                            "chunk_index": 1,
                        }
                    ],
                },
            ],
        },
    ],
}


def now_iso() -> str:
    return datetime.now().isoformat(timespec="seconds")


def make_result(test_id: str, requirement: str, name: str) -> dict[str, Any]:
    return {
        "id": test_id,
        "requirement": requirement,
        "name": name,
        "status": "NOT_RUN",
        "duration_s": None,
        "details": {},
        "error": "",
        "timestamp": now_iso(),
    }


def run_case(results: list[dict[str, Any]], test_id: str, requirement: str, name: str, fn: Callable[[], dict[str, Any] | None]) -> None:
    item = make_result(test_id, requirement, name)
    start = time.perf_counter()
    try:
        details = fn() or {}
        item["status"] = "PASS"
        item["details"] = details
    except AssertionError as exc:
        item["status"] = "FAIL"
        item["error"] = str(exc)
    except Exception as exc:  # noqa: BLE001 - test runner records unexpected failures.
        item["status"] = "ERROR"
        item["error"] = repr(exc)
    finally:
        item["duration_s"] = round(time.perf_counter() - start, 3)
        item["timestamp"] = now_iso()
        results.append(item)
        print(f"[{item['status']}] {test_id} {name} ({item['duration_s']}s)")
        if item["error"]:
            print(f"    {item['error']}")


def require_ok(resp: Any, *, expected_status: int = 200) -> dict[str, Any]:
    assert resp.status_code == expected_status, f"HTTP {resp.status_code}: {resp.text[:300]}"
    if resp.headers.get("content-type", "").startswith("application/json"):
        data = resp.json()
        assert data.get("ok", True) is not False, f"Response ok=false: {data}"
        return data
    return {}


def import_backend() -> Any:
    os.chdir(BACKEND_ROOT)
    sys.path.insert(0, str(BACKEND_ROOT))
    import http_server  # type: ignore

    return http_server


class RequestsClient:
    def __init__(self, base_url: str) -> None:
        self.base_url = base_url.rstrip("/")

    def get(self, path: str, **kwargs: Any) -> Any:
        import requests

        return requests.get(f"{self.base_url}{path}", timeout=kwargs.pop("timeout", 30), **kwargs)

    def post(self, path: str, **kwargs: Any) -> Any:
        import requests

        return requests.post(f"{self.base_url}{path}", timeout=kwargs.pop("timeout", 30), **kwargs)


def wait_for_server(base_url: str, timeout_s: int = 30) -> None:
    import requests

    deadline = time.time() + timeout_s
    last_error = ""
    while time.time() < deadline:
        try:
            resp = requests.post(f"{base_url}/api/ping", json={}, timeout=2)
            if resp.status_code == 200:
                return
            last_error = f"HTTP {resp.status_code}: {resp.text[:120]}"
        except Exception as exc:  # noqa: BLE001
            last_error = repr(exc)
        time.sleep(0.5)
    raise RuntimeError(f"server did not become ready: {last_error}")


def start_server(base_url: str) -> subprocess.Popen[str]:
    # Keep stdout/stderr captured so the test result is not flooded by uvicorn logs.
    proc = subprocess.Popen(
        [
            sys.executable,
            "-m",
            "uvicorn",
            "http_server:app",
            "--host",
            "127.0.0.1",
            "--port",
            base_url.rsplit(":", 1)[-1].replace("/", ""),
        ],
        cwd=BACKEND_ROOT,
        stdout=subprocess.PIPE,
        stderr=subprocess.STDOUT,
        text=True,
    )
    try:
        wait_for_server(base_url)
    except Exception:
        proc.terminate()
        try:
            out, _ = proc.communicate(timeout=5)
        except Exception:
            out = ""
        raise RuntimeError(f"failed to start uvicorn. Output:\n{out}")
    return proc


def run_local_tests(base_url: str | None = None) -> list[dict[str, Any]]:
    http_server = import_backend()
    from outline.evaluation import compute_quality
    from outline.schema import validate_outline

    server_proc: subprocess.Popen[str] | None = None
    if base_url is None:
        base_url = "http://127.0.0.1:18080"
        server_proc = start_server(base_url)
    client = RequestsClient(base_url)
    results: list[dict[str, Any]] = []
    conversation_id = f"acceptance-{uuid.uuid4().hex[:8]}"
    try:
        run_case(results, "TC-ENV-01", "Engineering", "Backend imports and FastAPI app loads", lambda: {
            "routes": len(http_server.app.routes),
            "env_file": str(http_server.ENV_FILE),
            "base_url": base_url,
        })

        run_case(results, "TC-API-01", "Engineering", "Ping endpoint", lambda: require_ok(client.post("/api/ping", json={})))
        run_case(results, "TC-API-02", "Engineering", "Runtime info endpoint", lambda: {
            **require_ok(client.get("/api/runtime-info")),
        })

        def schema_case() -> dict[str, Any]:
            ok, error = validate_outline(SAMPLE_OUTLINE)
            assert ok, error
            return {"chapter_count": len(SAMPLE_OUTLINE["chapters"])}

        run_case(results, "TC-SO-01", "Structured Outline", "Sample outline conforms to JSON schema", schema_case)

        def quality_case() -> dict[str, Any]:
            q = compute_quality(SAMPLE_OUTLINE, min_slides=6, max_slides=8)
            assert q.slide_count == 6, q
            assert q.chapter_count == 3, q
            assert q.overall_score_0_100 > 80, q
            return q.__dict__

        run_case(results, "TC-QA-02", "Quality Assurance", "Quality metrics are computed", quality_case)

        run_case(results, "TC-RAG-INDEX-01", "RAG", "List built RAG corpora", lambda: {
            "corpora": [c.get("id") for c in require_ok(client.get("/api/rag/corpora")).get("corpora", [])],
        })
        run_case(results, "TC-RAG-SRC-01", "RAG", "List RAG corpus sources", lambda: {
            "sources": [
                {"id": s.get("id"), "has_index": s.get("has_index"), "file_count": s.get("file_count")}
                for s in require_ok(client.get("/api/rag/corpus-sources")).get("sources", [])
            ],
        })

        def create_version() -> dict[str, Any]:
            data = require_ok(client.post("/api/outline/versions", json={
                "conversationId": conversation_id,
                "outline": SAMPLE_OUTLINE,
                "sourceType": "edited",
                "provider": "deepseek",
                "strategy": "baseline",
                "schemaMode": "on",
                "useRag": False,
                "summary": "Automated acceptance fixture",
            }))
            version = data.get("version") or {}
            assert version.get("versionId"), data
            create_version.version_id = version["versionId"]  # type: ignore[attr-defined]
            return {"conversation_id": conversation_id, "version": version}

        run_case(results, "TC-VER-01", "Engineering", "Create outline version", create_version)

        def list_versions() -> dict[str, Any]:
            data = require_ok(client.get(f"/api/outline/versions?conversationId={conversation_id}"))
            versions = data.get("versions") or []
            assert versions, data
            return {"count": len(versions), "latest": versions[0]}

        run_case(results, "TC-VER-02", "Engineering", "List outline versions", list_versions)

        def restore_version() -> dict[str, Any]:
            version_id = getattr(create_version, "version_id")
            data = require_ok(client.post(f"/api/outline/versions/{version_id}/restore", json={}))
            assert data.get("outline", {}).get("title") == SAMPLE_OUTLINE["title"], data
            assert data.get("version", {}).get("sourceType") == "restored", data
            return {"restored_version": data.get("version")}

        run_case(results, "TC-VER-03", "Engineering", "Restore outline version", restore_version)

        def export_case(fmt: str) -> dict[str, Any]:
            resp = client.post("/api/outline/export", json={
                "conversationId": conversation_id,
                "outline": SAMPLE_OUTLINE,
                "format": fmt,
                "report": {"source": "automated acceptance"},
            })
            assert resp.status_code == 200, f"HTTP {resp.status_code}: {resp.text[:300]}"
            assert len(resp.content) > 100, f"Export too small: {len(resp.content)} bytes"
            content_type = resp.headers.get("content-type", "")
            return {"format": fmt, "bytes": len(resp.content), "content_type": content_type}

        for fmt in ["markdown", "html", "json", "pptx"]:
            run_case(results, f"TC-EXP-{fmt.upper()}", "Engineering", f"Export {fmt}", lambda fmt=fmt: export_case(fmt))

        def extract_docx() -> dict[str, Any]:
            from docx import Document

            doc = Document()
            doc.add_paragraph("这是用于验收测试的 Word 文档。")
            doc.add_paragraph("主题：人工智能发展简史。")
            bio = io.BytesIO()
            doc.save(bio)
            bio.seek(0)
            data = require_ok(client.post("/api/extract/docx", files={
                "file": ("acceptance.docx", bio.getvalue(), "application/vnd.openxmlformats-officedocument.wordprocessingml.document")
            }))
            text = data.get("text", "")
            assert "人工智能发展简史" in text, data
            return {"chars": len(text)}

        run_case(results, "TC-IP-02-DOCX", "Input Processing", "Extract DOCX text", extract_docx)

        def extract_pptx() -> dict[str, Any]:
            from pptx import Presentation

            prs = Presentation()
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = "验收测试 PPT"
            slide.placeholders[1].text = "人工智能发展简史"
            bio = io.BytesIO()
            prs.save(bio)
            bio.seek(0)
            data = require_ok(client.post("/api/extract/pptx", files={
                "file": ("acceptance.pptx", bio.getvalue(), "application/vnd.openxmlformats-officedocument.presentationml.presentation")
            }))
            text = data.get("text", "")
            assert "人工智能发展简史" in text, data
            return {"chars": len(text)}

        run_case(results, "TC-IP-02-PPTX", "Input Processing", "Extract PPTX text", extract_pptx)

        return results
    finally:
        if server_proc is not None:
            server_proc.terminate()
            try:
                server_proc.wait(timeout=5)
            except subprocess.TimeoutExpired:
                server_proc.kill()


def write_results(results: list[dict[str, Any]], label: str) -> None:
    RESULTS_DIR.mkdir(exist_ok=True)
    stamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    json_path = RESULTS_DIR / f"{stamp}_{label}.json"
    md_path = RESULTS_DIR / f"{stamp}_{label}.md"

    summary: dict[str, int] = {}
    for item in results:
        summary[item["status"]] = summary.get(item["status"], 0) + 1

    payload = {
        "generated_at": now_iso(),
        "label": label,
        "summary": summary,
        "results": results,
    }
    json_path.write_text(json.dumps(payload, ensure_ascii=False, indent=2), encoding="utf-8")

    lines = [
        f"# Automated Acceptance Results ({label})",
        "",
        f"- Generated at: {payload['generated_at']}",
        f"- Summary: {summary}",
        "",
        "| ID | Requirement | Test | Status | Duration(s) | Notes |",
        "|---|---|---|---|---:|---|",
    ]
    for item in results:
        notes = item["error"] or json.dumps(item["details"], ensure_ascii=False)[:180]
        notes = notes.replace("|", "\\|").replace("\n", " ")
        lines.append(
            f"| {item['id']} | {item['requirement']} | {item['name']} | {item['status']} | {item['duration_s']} | {notes} |"
        )
    md_path.write_text("\n".join(lines) + "\n", encoding="utf-8")
    print(f"Wrote {json_path}")
    print(f"Wrote {md_path}")


def main() -> int:
    parser = argparse.ArgumentParser()
    parser.add_argument("--label", default="local")
    parser.add_argument("--base-url", default=None)
    args = parser.parse_args()
    results = run_local_tests(args.base_url)
    write_results(results, args.label)
    return 0 if all(r["status"] == "PASS" for r in results) else 1


if __name__ == "__main__":
    raise SystemExit(main())
